from flask import Flask, request, send_file, jsonify
from pptx import Presentation
from pptx.util import Inches
import requests
import tempfile
import os
import shutil
import zipfile
import subprocess

app = Flask(__name__)

@app.route("/generate", methods=["POST"])
def generate_ppt():
    try:
        data = request.get_json()
        pptx_url = data.get("pptx_url")
        qr_url = data.get("qr_url")

        if not pptx_url or not qr_url:
            return jsonify({"error": "Missing pptx_url or qr_url"}), 400

        with tempfile.TemporaryDirectory() as tmpdir:
            # Download the pptx and QR
            pptx_path = os.path.join(tmpdir, "template.pptx")
            qr_path = os.path.join(tmpdir, "qr.png")

            with open(pptx_path, "wb") as f:
                f.write(requests.get(pptx_url).content)

            with open(qr_path, "wb") as f:
                f.write(requests.get(qr_url).content)

            # Open presentation and insert QR
            prs = Presentation(pptx_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame and "{{QR}}" in shape.text:
                        left = shape.left
                        top = shape.top
                        height = shape.height
                        width = shape.width
                        slide.shapes._spTree.remove(shape._element)
                        slide.shapes.add_picture(qr_path, left, top, width, height)
            pptx_output_path = os.path.join(tmpdir, "Kit de Bienvenida - output.pptx")
            prs.save(pptx_output_path)

            # Convert to PDF using LibreOffice
            pdf_output_path = pptx_output_path.replace(".pptx", ".pdf")
            subprocess.run([
                "libreoffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", tmpdir,
                pptx_output_path
            ], check=True)

            # Create zip with both files
            zip_path = os.path.join(tmpdir, "output_files.zip")
            with zipfile.ZipFile(zip_path, 'w') as zipf:
                zipf.write(pptx_output_path, arcname="Kit de Bienvenida - output.pptx")
                zipf.write(pdf_output_path, arcname="Kit de Bienvenida - output.pdf")

            return send_file(zip_path, as_attachment=True, download_name="output_files.zip")

    except Exception as e:
        return jsonify({"error": str(e)}), 500

